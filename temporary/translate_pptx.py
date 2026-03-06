from pptx import Presentation

prs = Presentation(r'C:\Users\omar.ola\NotificationAPI\temporary\Notification_API_Microservices_Overview.pptx')
VT = '\x0b'  # vertical tab used in pptx line breaks

translations = {
    # Slide 1
    "Spec-Driven Development vs Vibe Coding": "Desarrollo Guiado por Especificaciones vs Vibe Coding",
    "Two AI-assisted development approaches \u2014 same tools, fundamentally different outcomes": "Dos enfoques de desarrollo asistido por IA \u2014 mismas herramientas, resultados fundamentalmente diferentes",
    '"Prompt and see what happens"': '"Escribe un prompt y mira qu\u00e9 pasa"',
    "Loose prompts, no upfront design, iterate by feel.Fast for prototypes \u2014 but inconsistent quality,no test strategy, hard to maintain past ~1K LoC.Best for: hackathons, MVPs, throwaway projects.":
        "Prompts vagos, sin dise\u00f1o previo, iteraci\u00f3n por intuici\u00f3n. R\u00e1pido para prototipos \u2014 pero calidad inconsistente, sin estrategia de pruebas, dif\u00edcil de mantener m\u00e1s all\u00e1 de ~1K LoC. Ideal para: hackathons, MVPs, proyectos desechables.",
    "Spec-Driven Development": "Desarrollo Guiado por Especificaciones",
    '"Design first, then instruct the AI precisely"': '"Primero dise\u00f1a, luego instruye a la IA con precisi\u00f3n"',
    "Detailed design docs, DB schemas, API contracts \u2014all defined before writing line 1 of code.AI executes against precise specifications.Best for: production systems, teams, long-lived projects.":
        "Documentos de dise\u00f1o detallados, esquemas de BD, contratos de API \u2014 todo definido antes de escribir la l\u00ednea 1 de c\u00f3digo. La IA ejecuta contra especificaciones precisas. Ideal para: sistemas en producci\u00f3n, equipos, proyectos de larga duraci\u00f3n.",
    "Why Spec-Driven Development Matters": "Por qu\u00e9 Importa el Desarrollo Guiado por Especificaciones",
    "AI is only as good as its instructions.": "La IA es tan buena como sus instrucciones.",
    "Vague prompts produce vague code. Detailed specs give the AI exact constraints \u2014 data models, error codes, API contracts, queue topologies \u2014 so it generates production-grade code on the first pass, not the fifth.":
        "Prompts vagos producen c\u00f3digo vago. Las especificaciones detalladas le dan a la IA restricciones exactas \u2014 modelos de datos, c\u00f3digos de error, contratos de API, topolog\u00edas de colas \u2014 para que genere c\u00f3digo de nivel producci\u00f3n en el primer intento, no en el quinto.",
    "Specs are the single source of truth for the entire team.": "Las especificaciones son la \u00fanica fuente de verdad para todo el equipo.",
    "When every service has a design doc, DB script, and endpoint reference, any developer (or AI) can pick up any service and know exactly what it should do. No tribal knowledge, no guessing.":
        "Cuando cada servicio tiene un documento de dise\u00f1o, script de BD y referencia de endpoints, cualquier desarrollador (o IA) puede tomar cualquier servicio y saber exactamente qu\u00e9 debe hacer. Sin conocimiento tribal, sin adivinanzas.",
    "You can scale without chaos.": "Puedes escalar sin caos.",
    "This project: 14 microservices, 45K source LoC, 59K test LoC, 2,375+ unit tests \u2014 designed and built in ~15 days. That velocity is only possible because specs eliminated ambiguity and rework.":
        "Este proyecto: 14 microservicios, 45K LoC fuente, 59K LoC de pruebas, 2,375+ pruebas unitarias \u2014 dise\u00f1ado y construido en ~15 d\u00edas. Esa velocidad solo es posible porque las especificaciones eliminaron la ambig\u00fcedad y el retrabajo.",
    "Specs survive the code.": "Las especificaciones sobreviven al c\u00f3digo.",
    "Code gets refactored, rewritten, or replaced. Specs document the WHY and the WHAT \u2014 they remain valuable long after the current implementation changes. They are the institutional memory of the system.":
        "El c\u00f3digo se refactoriza, reescribe o reemplaza. Las especificaciones documentan el POR QU\u00c9 y el QU\u00c9 \u2014 siguen siendo valiosas mucho despu\u00e9s de que cambie la implementaci\u00f3n actual. Son la memoria institucional del sistema.",

    # Slide 2
    "Microservices Platform Overview": "Visi\u00f3n General de la Plataforma de Microservicios",
    "Unified notification platform for eCommerce - consolidates fragmented notificationsfrom multiple source systems into a single event-driven microservices architecture.":
        "Plataforma unificada de notificaciones para eCommerce - consolida notificaciones fragmentadas de m\u00faltiples sistemas fuente en una \u00fanica arquitectura de microservicios basada en eventos.",
    "Unified notification platform for eCommerce - consolidates fragmented notifications\x0bfrom multiple source systems into a single event-driven microservices architecture.":
        "Plataforma unificada de notificaciones para eCommerce - consolida notificaciones fragmentadas\x0bde m\u00faltiples sistemas fuente en una \u00fanica arquitectura de microservicios basada en eventos.",
    "Total Services": "Servicios Totales",
    "Implemented": "Implementados",
    "Pending": "Pendientes",
    "Lines of Code": "L\u00edneas de C\u00f3digo",
    "Design: Feb 19 - Feb 28, 2026  |  Development: Feb 23 - Mar 5, 2026": "Dise\u00f1o: Feb 19 - Feb 28, 2026  |  Desarrollo: Feb 23 - Mar 5, 2026",

    # Slide 3
    "Notification Flow \u2014 End to End": "Flujo de Notificaci\u00f3n \u2014 De Extremo a Extremo",
    "Event-driven pipeline: Source Systems  >>  Ingestion  >>  Engine  >>  Rendering  >>  Routing  >>  Delivery":
        "Pipeline basado en eventos: Sistemas Fuente  >>  Ingesti\u00f3n  >>  Motor  >>  Renderizado  >>  Enrutamiento  >>  Entrega",
    "Core Pipeline": "Pipeline Principal",
    "Message Broker": "Broker de Mensajes",
    "Support Services": "Servicios de Soporte",
    "Admin / Auth": "Admin / Auth",
    "Sources / Ingest": "Fuentes / Ingesti\u00f3n",
    "External Providers": "Proveedores Externos",
    "Normalize & Validate": "Normalizar y Validar",
    "Rules & Orchestrate": "Reglas y Orquestar",
    "Route & Deliver": "Enrutar y Entregar",
    "Render (Handlebars)": "Renderizar (Handlebars)",
    "XLSX Processing": "Procesamiento XLSX",
    "SMTP Ingest": "Ingesti\u00f3n SMTP",
    "Tracking & Analytics": "Seguimiento y An\u00e1lisis",
    "Config & Rules": "Configuraci\u00f3n y Reglas",
    "JWT & Users": "JWT y Usuarios",
    # Slide 3 diagram boxes (with \x0b vertical tabs)
    "Event Ingestion\x0bService\x0b:3151\x0bNormalize & Validate": "Event Ingestion\x0bService\x0b:3151\x0bNormalizar y Validar",
    "Notification\x0bEngine\x0b:3152\x0bRules & Orchestrate": "Notification\x0bEngine\x0b:3152\x0bReglas y Orquestar",
    "Channel Router\x0bService\x0b:3154\x0bRoute & Deliver": "Channel Router\x0bService\x0b:3154\x0bEnrutar y Entregar",
    "Template Service\x0b:3153\x0bRender (Handlebars)": "Template Service\x0b:3153\x0bRenderizar (Handlebars)",
    "Bulk Upload\x0bService :3158\x0bXLSX Processing": "Bulk Upload\x0bService :3158\x0bProcesamiento XLSX",
    "Email Ingest\x0bService :3157/2525\x0bSMTP Ingest": "Email Ingest\x0bService :3157/2525\x0bIngesti\u00f3n SMTP",
    "Audit Service\x0b:3156\x0bTracking & Analytics": "Audit Service\x0b:3156\x0bSeguimiento y An\u00e1lisis",
    "Admin Service\x0b:3155\x0bConfig & Rules": "Admin Service\x0b:3155\x0bConfiguraci\u00f3n y Reglas",
    "Auth RBAC\x0bBackend :3160\x0bJWT & Users": "Auth RBAC\x0bBackend :3160\x0bJWT y Usuarios",
    "Source\x0bSystems\x0bOMS\x0bMagento\x0bMirakl\x0bChat\x0bManual": "Sistemas\x0bFuente\x0bOMS\x0bMagento\x0bMirakl\x0bChat\x0bManual",
    "Provider\x0bAdapters\x0bMailgun\x0bWhatsApp\x0bBraze\x0bAWS SES": "Adaptadores\x0bde Proveedor\x0bMailgun\x0bWhatsApp\x0bBraze\x0bAWS SES",
    "RabbitMQ\x0b6 Exchanges\x0b23 Queues\x0b24 Bindings": "RabbitMQ\x0b6 Exchanges\x0b23 Colas\x0b24 Bindings",

    # Slide 4
    "Development Timeline": "Cronograma de Desarrollo",
    "Design: Feb 19-28  |  Implementation: Feb 23 - Mar 5  |  Total: ~15 days (design + dev)":
        "Dise\u00f1o: Feb 19-28  |  Implementaci\u00f3n: Feb 23 - Mar 5  |  Total: ~15 d\u00edas (dise\u00f1o + desarrollo)",
    "Design Phase (all services)": "Fase de Dise\u00f1o (todos los servicios)",
    "Services were developed in parallel streams with overlapping timelines. Design docs preceded implementation by 2-4 days per service.":
        "Los servicios se desarrollaron en flujos paralelos con cronogramas superpuestos. Los documentos de dise\u00f1o precedieron a la implementaci\u00f3n por 2-4 d\u00edas por servicio.",

    # Slide 5
    "Implemented Microservices": "Microservicios Implementados",
    "8 services fully or partially implemented with production-ready code and tests":
        "8 servicios total o parcialmente implementados con c\u00f3digo listo para producci\u00f3n y pruebas",
    "Service": "Servicio",
    "Port": "Puerto",
    "Endpoints": "Endpoints",
    "Source LoC": "LoC Fuente",
    "Test LoC": "LoC Pruebas",
    "Unit Tests": "Pruebas Unitarias",
    "Dev Time": "Tiempo Dev",
    "Status": "Estado",
    "Complete (Step 4)": "Completo (Paso 4)",
    "Complete (Phase 7)": "Completo (Fase 7)",
    "Complete (Phase 5)": "Completo (Fase 5)",
    "Complete (Phase 4)": "Completo (Fase 4)",
    "Partial (MG+WA)": "Parcial (MG+WA)",
    "Phase 13 (87%)": "Fase 13 (87%)",

    # Slide 6
    "Services Yet to Be Implemented": "Servicios Pendientes de Implementaci\u00f3n",
    "6 services in scaffolding/design phase \u2014 design docs complete, code pending":
        "6 servicios en fase de andamiaje/dise\u00f1o \u2014 documentos de dise\u00f1o completos, c\u00f3digo pendiente",
    "Planned Endpoints": "Endpoints Planificados",
    "Design Doc": "Doc. Dise\u00f1o",
    "Purpose": "Prop\u00f3sito",
    "Estimated Complexity": "Complejidad Estimada",
    "Backoffice config management (rules, mappings, channels)": "Gesti\u00f3n de configuraci\u00f3n backoffice (reglas, mapeos, canales)",
    "High (39 endpoints, CRUD+)": "Alta (39 endpoints, CRUD+)",
    "SMTP ingest, email parsing, event generation": "Ingesti\u00f3n SMTP, an\u00e1lisis de emails, generaci\u00f3n de eventos",
    "Medium (SMTP + parsing)": "Media (SMTP + an\u00e1lisis)",
    "Multi-app auth, JWT RS256, user & role mgmt": "Auth multi-app, JWT RS256, gesti\u00f3n de usuarios y roles",
    "High (auth + RBAC)": "Alta (auth + RBAC)",
    "Next.js admin UI for auth/RBAC management": "UI admin Next.js para gesti\u00f3n de auth/RBAC",
    "Medium (6 pages)": "Media (6 p\u00e1ginas)",
    "Login portal & app launcher": "Portal de login y lanzador de apps",
    "Low (4 pages)": "Baja (4 p\u00e1ginas)",
    "BFF / API Gateway (DEPRECATED)": "BFF / API Gateway (DEPRECADO)",
    "N/A \u2014 Deprecated": "N/A \u2014 Deprecado",
    "Note: notification-gateway is deprecated. Its responsibilities were redistributed: auth to auth-rbac-service, RBAC to per-service JWT validation, rate limiting to infrastructure proxy.":
        "Nota: notification-gateway est\u00e1 deprecado. Sus responsabilidades fueron redistribuidas: auth a auth-rbac-service, RBAC a validaci\u00f3n JWT por servicio, limitaci\u00f3n de tasa a proxy de infraestructura.",

    # Slide 7
    "Complexity Analysis (Lines of Code)": "An\u00e1lisis de Complejidad (L\u00edneas de C\u00f3digo)",
    "Largest Service": "Servicio M\u00e1s Grande",
    "notification-admin-ui15,791 LoC (Next.js frontend)": "notification-admin-ui\n15,791 LoC (Frontend Next.js)",
    "notification-admin-ui\x0b15,791 LoC (Next.js frontend)": "notification-admin-ui\x0b15,791 LoC (Frontend Next.js)",
    "Most Complex Backend": "Backend M\u00e1s Complejo",
    "notification-engine-service6,020 src + 11,575 test LoC504 unit tests across 51 suites":
        "notification-engine-service\n6,020 src + 11,575 test LoC\n504 pruebas unitarias en 51 suites",
    "notification-engine-service\x0b6,020 src + 11,575 test LoC\x0b504 unit tests across 51 suites":
        "notification-engine-service\x0b6,020 src + 11,575 test LoC\x0b504 pruebas unitarias en 51 suites",
    "Best Test Ratio": "Mejor Ratio de Pruebas",
    "bulk-upload-service2.26x test-to-source ratio": "bulk-upload-service\n2.26x ratio pruebas-a-fuente",
    "bulk-upload-service\x0b2.26x test-to-source ratio": "bulk-upload-service\x0b2.26x ratio pruebas-a-fuente",
    "Total Codebase": "C\u00f3digo Total",
    "~45K source + ~59K test= ~104K total lines of code": "~45K fuente + ~59K pruebas\n= ~104K l\u00edneas de c\u00f3digo totales",
    "~45K source + ~59K test\x0b= ~104K total lines of code": "~45K fuente + ~59K pruebas\x0b= ~104K l\u00edneas de c\u00f3digo totales",

    # Slide 8
    "Endpoint Coverage Summary": "Resumen de Cobertura de Endpoints",
    "Implemented Endpoints": "Endpoints Implementados",
    "of 266 total (45%)": "de 266 totales (45%)",
    "Pending Endpoints": "Endpoints Pendientes",
    "across 5 services": "en 5 servicios",
    "Deprecated Endpoints": "Endpoints Deprecados",
    "Endpoints per Service (Implemented vs Planned)": "Endpoints por Servicio (Implementados vs Planificados)",

    # Slide 9
    "Key Takeaways": "Conclusiones Clave",
    "Architecture": "Arquitectura",
    "14 microservices designed, 8 implemented. Event-driven with RabbitMQ (6 exchanges, 23 queues). PostgreSQL with schema-per-service pattern. NestJS + TypeScript backend, Next.js frontend.":
        "14 microservicios dise\u00f1ados, 8 implementados. Basado en eventos con RabbitMQ (6 exchanges, 23 colas). PostgreSQL con patr\u00f3n schema-por-servicio. Backend NestJS + TypeScript, frontend Next.js.",
    "Code Quality": "Calidad de C\u00f3digo",
    "~45K lines of production source code, ~59K lines of test code. Test-to-source ratio averages 1.3x-2.3x across services. 2,375+ unit tests across 243 suites, plus E2E test coverage.":
        "~45K l\u00edneas de c\u00f3digo fuente de producci\u00f3n, ~59K l\u00edneas de c\u00f3digo de pruebas. El ratio pruebas-a-fuente promedia 1.3x-2.3x entre servicios. 2,375+ pruebas unitarias en 243 suites, m\u00e1s cobertura de pruebas E2E.",
    "Development Velocity": "Velocidad de Desarrollo",
    "Full platform designed in ~10 days (Feb 19-28). 8 services implemented in ~10 days (Feb 23 - Mar 5). Parallel development streams enabled rapid delivery.":
        "Plataforma completa dise\u00f1ada en ~10 d\u00edas (Feb 19-28). 8 servicios implementados en ~10 d\u00edas (Feb 23 - Mar 5). Flujos de desarrollo paralelos permitieron una entrega r\u00e1pida.",
    "Next Steps": "Pr\u00f3ximos Pasos",
    "Priority: admin-service (39 endpoints, config management hub), auth-rbac-service-backend (19 endpoints, platform security). Then: auth-rbac-frontend, ecommerce-backoffice, email-ingest-service. Remaining provider adapters: Braze, AWS SES.":
        "Prioridad: admin-service (39 endpoints, hub de gesti\u00f3n de configuraci\u00f3n), auth-rbac-service-backend (19 endpoints, seguridad de plataforma). Luego: auth-rbac-frontend, ecommerce-backoffice, email-ingest-service. Adaptadores de proveedor restantes: Braze, AWS SES.",
}


def translate_text(text):
    stripped = text.strip()
    if stripped in translations:
        return translations[stripped]
    return text


def translate_runs(paragraph):
    full_text = paragraph.text.strip()
    if not full_text:
        return

    translated = translate_text(full_text)
    if translated == full_text:
        return

    if len(paragraph.runs) <= 1:
        if paragraph.runs:
            paragraph.runs[0].text = translated
        return

    # Multiple runs: put all text in first run, clear the rest
    paragraph.runs[0].text = translated
    for run in paragraph.runs[1:]:
        run.text = ""


for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                translate_runs(para)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        translate_runs(para)

output_path = r'C:\Users\omar.ola\NotificationAPI\temporary\Notification_API_Microservices_Overview_ES.pptx'
prs.save(output_path)
print(f"Saved to: {output_path}")
