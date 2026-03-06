"""Generate PowerPoint presentation for Notification API Microservices Overview."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_GREEN = RGBColor(0x00, 0xC9, 0x7B)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xE8, 0x4D, 0x4D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x40)
TABLE_HEADER = RGBColor(0x00, 0x64, 0x8C)
TABLE_ROW1 = RGBColor(0x1E, 0x1E, 0x35)
TABLE_ROW2 = RGBColor(0x28, 0x28, 0x45)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_styled_table(slide, left, top, width, rows, cols, data, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                          Inches(width), Inches(0.4 * rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(data[row_idx][col_idx])

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Segoe UI"
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = LIGHT_GRAY
                    paragraph.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER

            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER
            elif row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW1
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW2

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


# ─── SLIDE 1: Title ───
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# Accent line
add_shape_rect(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)

add_text_box(slide, 1, 1.5, 11, 1.2, "Notification API", font_size=48, color=WHITE, bold=True)
add_text_box(slide, 1, 2.7, 11, 0.8, "Microservices Platform Overview", font_size=28, color=ACCENT_BLUE)
add_text_box(slide, 1, 3.8, 11, 0.6,
             "Unified notification platform for eCommerce - consolidates fragmented notifications\n"
             "from multiple source systems into a single event-driven microservices architecture.",
             font_size=16, color=LIGHT_GRAY)

# Key stats cards
stats = [
    ("14", "Total Services"),
    ("8", "Implemented"),
    ("6", "Pending"),
    ("~104K", "Lines of Code"),
]
for i, (num, label) in enumerate(stats):
    x = 1.5 + i * 2.7
    add_shape_rect(slide, x, 5.2, 2.2, 1.3, CARD_BG, ACCENT_BLUE)
    add_text_box(slide, x, 5.3, 2.2, 0.7, num, font_size=32, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, 5.9, 2.2, 0.5, label, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1, 6.9, 11, 0.4, "Design: Feb 19 - Feb 28, 2026  |  Development: Feb 23 - Mar 5, 2026",
             font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 2: Implemented Services Table ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_rect(slide, 0, 0, 13.333, 0.08, ACCENT_GREEN)
add_text_box(slide, 0.5, 0.3, 12, 0.7, "Implemented Microservices", font_size=32, color=WHITE, bold=True)
add_text_box(slide, 0.5, 0.9, 12, 0.4, "8 services fully or partially implemented with production-ready code and tests",
             font_size=14, color=LIGHT_GRAY)

data = [
    ["Service", "Port", "Endpoints", "Source LoC", "Test LoC", "Unit Tests", "Dev Time", "Status"],
    ["event-ingestion-service", "3151", "11", "3,725", "6,478", "264 / 33 suites", "~7 days", "Complete (Step 4)"],
    ["notification-engine-service", "3152", "35", "6,020", "11,575", "504 / 51 suites", "~7 days", "Complete (Phase 7)"],
    ["template-service", "3153", "10", "2,816", "6,336", "250 / 21 suites", "~6 days", "Complete (Phase 5)"],
    ["channel-router-service", "3154", "12 (v2)", "4,943", "9,633", "432 / 46 suites", "~5 days", "Complete (Phase 4)"],
    ["provider-adapters", "—", "8/20", "4,399", "6,461", "204 / 28 suites", "~5 days", "Partial (MG+WA)"],
    ["audit-service", "3156", "12", "4,098", "8,619", "338 / 41 suites", "~3 days", "Complete (Phase 4)"],
    ["bulk-upload-service", "3158", "10", "4,339", "9,797", "387 / 23 suites", "~3 days", "Complete (Phase 4)"],
    ["notification-admin-ui", "3159", "21/24", "15,791", "—", "—", "~5 days", "Phase 13 (87%)"],
]

add_styled_table(slide, 0.3, 1.5, 12.7, len(data), 8, data,
                 col_widths=[2.5, 0.6, 1.0, 1.1, 1.1, 1.8, 1.0, 1.8])


# ─── SLIDE 3: Pending Services ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_rect(slide, 0, 0, 13.333, 0.08, ACCENT_ORANGE)
add_text_box(slide, 0.5, 0.3, 12, 0.7, "Services Yet to Be Implemented", font_size=32, color=WHITE, bold=True)
add_text_box(slide, 0.5, 0.9, 12, 0.4, "6 services in scaffolding/design phase — design docs complete, code pending",
             font_size=14, color=LIGHT_GRAY)

data2 = [
    ["Service", "Port", "Planned Endpoints", "Design Doc", "Purpose", "Estimated Complexity"],
    ["admin-service", "3155", "39", "14-admin-service.md", "Backoffice config management (rules, mappings, channels)", "High (39 endpoints, CRUD+)"],
    ["email-ingest-service", "3157", "6", "—", "SMTP ingest, email parsing, event generation", "Medium (SMTP + parsing)"],
    ["auth-rbac-service-backend", "3160", "19", "19-auth-rbac-service.md", "Multi-app auth, JWT RS256, user & role mgmt", "High (auth + RBAC)"],
    ["auth-rbac-service-frontend", "3161", "6", "20-auth-rbac-frontend.md", "Next.js admin UI for auth/RBAC management", "Medium (6 pages)"],
    ["ecommerce-backoffice", "3162", "4", "21-ecommerce-backoffice.md", "Login portal & app launcher", "Low (4 pages)"],
    ["notification-gateway", "3150", "49", "13-notification-gateway.md", "BFF / API Gateway (DEPRECATED)", "N/A — Deprecated"],
]

add_styled_table(slide, 0.3, 1.5, 12.7, len(data2), 6, data2,
                 col_widths=[2.5, 0.7, 1.5, 2.2, 3.5, 2.3])

# Note
add_text_box(slide, 0.5, 5.5, 12, 0.8,
             "Note: notification-gateway is deprecated. Its responsibilities were redistributed: "
             "auth to auth-rbac-service, RBAC to per-service JWT validation, rate limiting to infrastructure proxy.",
             font_size=11, color=MED_GRAY)


# ─── SLIDE 4: Complexity Analysis ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_rect(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)
add_text_box(slide, 0.5, 0.3, 12, 0.7, "Complexity Analysis (Lines of Code)", font_size=32, color=WHITE, bold=True)

# Bar chart data
chart_data = CategoryChartData()
chart_data.categories = [
    'event-ingestion', 'notif-engine', 'template', 'channel-router',
    'provider-adapters', 'audit', 'bulk-upload', 'admin-ui'
]
chart_data.add_series('Source Code', (3725, 6020, 2816, 4943, 4399, 4098, 4339, 15791))
chart_data.add_series('Test Code', (6478, 11575, 6336, 9633, 6461, 8619, 9797, 0))

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.5), Inches(1.3), Inches(8), Inches(5.5),
    chart_data
).chart

chart.has_legend = True
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(11)
chart.legend.font.color.rgb = LIGHT_GRAY

plot = chart.plots[0]
plot.gap_width = 80
series_src = plot.series[0]
series_test = plot.series[1]
series_src.format.fill.solid()
series_src.format.fill.fore_color.rgb = ACCENT_BLUE
series_test.format.fill.solid()
series_test.format.fill.fore_color.rgb = ACCENT_GREEN

chart.category_axis.tick_labels.font.size = Pt(9)
chart.category_axis.tick_labels.font.color.rgb = LIGHT_GRAY
chart.value_axis.tick_labels.font.size = Pt(9)
chart.value_axis.tick_labels.font.color.rgb = LIGHT_GRAY

# Insights panel
add_shape_rect(slide, 8.8, 1.3, 4.2, 5.5, CARD_BG, ACCENT_BLUE)
insights = [
    ("Largest Service", "notification-admin-ui\n15,791 LoC (Next.js frontend)"),
    ("Most Complex Backend", "notification-engine-service\n6,020 src + 11,575 test LoC\n504 unit tests across 51 suites"),
    ("Best Test Ratio", "bulk-upload-service\n2.26x test-to-source ratio"),
    ("Total Codebase", "~45K source + ~59K test\n= ~104K total lines of code"),
]
y = 1.5
for title, desc in insights:
    add_text_box(slide, 9.0, y, 3.8, 0.3, title, font_size=13, color=ACCENT_GREEN, bold=True)
    add_text_box(slide, 9.0, y + 0.3, 3.8, 0.7, desc, font_size=10, color=LIGHT_GRAY)
    y += 1.25


# ─── SLIDE 5: Development Timeline ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_rect(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)
add_text_box(slide, 0.5, 0.3, 12, 0.7, "Development Timeline", font_size=32, color=WHITE, bold=True)
add_text_box(slide, 0.5, 0.9, 12, 0.4,
             "Design: Feb 19-28  |  Implementation: Feb 23 - Mar 5  |  Total: ~15 days (design + dev)",
             font_size=14, color=LIGHT_GRAY)

# Timeline visualization as horizontal bars
services_timeline = [
    ("Design Phase (all services)", "Feb 19", "Feb 28", 0, 9, ACCENT_ORANGE),
    ("event-ingestion-service", "Feb 23", "Mar 1", 4, 7, ACCENT_BLUE),
    ("notification-engine-service", "Feb 24", "Mar 2", 5, 7, ACCENT_BLUE),
    ("template-service", "Feb 25", "Mar 2", 6, 6, ACCENT_BLUE),
    ("channel-router-service", "Feb 26", "Mar 2", 7, 5, ACCENT_BLUE),
    ("provider-adapters", "Feb 26", "Mar 2", 7, 5, ACCENT_GREEN),
    ("audit-service", "Feb 27", "Mar 1", 8, 3, ACCENT_GREEN),
    ("bulk-upload-service", "Feb 27", "Mar 1", 8, 3, ACCENT_GREEN),
    ("notification-admin-ui", "Feb 27", "Mar 1", 8, 3, ACCENT_GREEN),
]

# Date labels
bar_left = 3.5
bar_max_w = 9.0
total_days = 15  # Feb 19 to Mar 5
dates_labels = ["Feb 19", "Feb 21", "Feb 23", "Feb 25", "Feb 27", "Mar 1", "Mar 3", "Mar 5"]
for i, d in enumerate(dates_labels):
    x = bar_left + (i * (bar_max_w / (len(dates_labels) - 1)))
    add_text_box(slide, x - 0.3, 1.3, 0.8, 0.3, d, font_size=8, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

y_start = 1.8
for i, (name, start, end, offset_days, duration, color) in enumerate(services_timeline):
    y = y_start + i * 0.55
    add_text_box(slide, 0.3, y, 3.0, 0.4, name, font_size=11, color=WHITE, bold=(i == 0))

    x = bar_left + (offset_days / total_days) * bar_max_w
    w = max((duration / total_days) * bar_max_w, 0.3)
    add_shape_rect(slide, x, y + 0.05, w, 0.3, color)

    add_text_box(slide, x + 0.1, y + 0.05, w - 0.2, 0.3,
                 f"{start} - {end}", font_size=8, color=WHITE, alignment=PP_ALIGN.CENTER)

# Bottom note
add_text_box(slide, 0.5, 6.8, 12, 0.4,
             "Services were developed in parallel streams with overlapping timelines. "
             "Design docs preceded implementation by 2-4 days per service.",
             font_size=11, color=MED_GRAY)


# ─── SLIDE 6: Architecture & Endpoint Summary ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_rect(slide, 0, 0, 13.333, 0.08, ACCENT_GREEN)
add_text_box(slide, 0.5, 0.3, 12, 0.7, "Endpoint Coverage Summary", font_size=32, color=WHITE, bold=True)

# Donut-style summary cards
cards = [
    ("Implemented Endpoints", "119", "of 266 total (45%)", ACCENT_GREEN),
    ("Pending Endpoints", "98", "across 5 services", ACCENT_ORANGE),
    ("Deprecated Endpoints", "49", "notification-gateway", ACCENT_RED),
]
for i, (title, num, desc, color) in enumerate(cards):
    x = 0.5 + i * 4.3
    add_shape_rect(slide, x, 1.2, 3.8, 1.8, CARD_BG, color)
    add_text_box(slide, x, 1.3, 3.8, 0.4, title, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, 1.7, 3.8, 0.7, num, font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, 2.4, 3.8, 0.4, desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Endpoint breakdown by service
add_text_box(slide, 0.5, 3.3, 12, 0.5, "Endpoints per Service (Implemented vs Planned)",
             font_size=18, color=WHITE, bold=True)

chart_data2 = CategoryChartData()
chart_data2.categories = [
    'event-\ningestion', 'notif-\nengine', 'template', 'channel-\nrouter v2',
    'provider-\nadapters', 'audit', 'bulk-\nupload', 'admin-\nui',
    'admin-\nservice', 'email-\ningest', 'auth-\nbackend', 'auth-\nfrontend', 'ecomm-\nbackoffice'
]
chart_data2.add_series('Done', (11, 35, 10, 12, 8, 12, 10, 21, 0, 0, 0, 0, 0))
chart_data2.add_series('Pending', (0, 0, 0, 0, 12, 0, 0, 3, 39, 6, 19, 6, 4))

chart2 = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_STACKED,
    Inches(0.5), Inches(3.8), Inches(12.3), Inches(3.3),
    chart_data2
).chart

chart2.has_legend = True
chart2.legend.include_in_layout = False
chart2.legend.font.size = Pt(11)
chart2.legend.font.color.rgb = LIGHT_GRAY

plot2 = chart2.plots[0]
plot2.gap_width = 60
s1 = plot2.series[0]
s2 = plot2.series[1]
s1.format.fill.solid()
s1.format.fill.fore_color.rgb = ACCENT_GREEN
s2.format.fill.solid()
s2.format.fill.fore_color.rgb = ACCENT_ORANGE

chart2.category_axis.tick_labels.font.size = Pt(8)
chart2.category_axis.tick_labels.font.color.rgb = LIGHT_GRAY
chart2.value_axis.tick_labels.font.size = Pt(9)
chart2.value_axis.tick_labels.font.color.rgb = LIGHT_GRAY


# ─── SLIDE 7: Notification Flow Diagram ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_rect(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)
add_text_box(slide, 0.5, 0.2, 12, 0.6, "Notification Flow — End to End", font_size=32, color=WHITE, bold=True)
add_text_box(slide, 0.5, 0.75, 12, 0.3,
             "Event-driven pipeline: Source Systems  >>  Ingestion  >>  Engine  >>  Rendering  >>  Routing  >>  Delivery",
             font_size=12, color=LIGHT_GRAY)

# --- Helper: draw a box with centered label ---
from pptx.enum.shapes import MSO_SHAPE


def add_box(slide, left, top, width, height, label, fill_color, border_color, font_size=10, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = "Segoe UI"
    p.font.bold = True
    # vertical center
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT_BLUE):
    """Draw connector line from (x1,y1) to (x2,y2) in inches."""
    from pptx.enum.shapes import MSO_SHAPE
    if abs(x2 - x1) > abs(y2 - y1):
        # horizontal arrow
        left = min(x1, x2)
        top = min(y1, y2) - 0.02
        w = abs(x2 - x1)
        h = 0.04
    else:
        # vertical arrow
        left = min(x1, x2) - 0.02
        top = min(y1, y2)
        w = 0.04
        h = abs(y2 - y1)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_arrow_head(slide, cx, cy, direction, color=ACCENT_BLUE):
    """Small triangle pointing right or down."""
    from pptx.enum.shapes import MSO_SHAPE
    if direction == "right":
        shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                        Inches(cx), Inches(cy - 0.06), Inches(0.12), Inches(0.12))
        shape.rotation = 90.0
    elif direction == "down":
        shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                        Inches(cx - 0.06), Inches(cy), Inches(0.12), Inches(0.12))
        shape.rotation = 180.0
    elif direction == "left":
        shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                        Inches(cx - 0.12), Inches(cy - 0.06), Inches(0.12), Inches(0.12))
        shape.rotation = 270.0
    elif direction == "up":
        shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                        Inches(cx - 0.06), Inches(cy - 0.12), Inches(0.12), Inches(0.12))
        shape.rotation = 0.0
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# Color definitions for box types
SRC_COLOR = RGBColor(0x3A, 0x3A, 0x5C)       # source systems
SRC_BORDER = RGBColor(0x88, 0x88, 0xBB)
CORE_COLOR = RGBColor(0x00, 0x50, 0x80)       # core pipeline
CORE_BORDER = ACCENT_BLUE
MQ_COLOR = RGBColor(0xFF, 0x6B, 0x00)         # RabbitMQ
MQ_BORDER = RGBColor(0xFF, 0x9E, 0x50)
SUPPORT_COLOR = RGBColor(0x1A, 0x6B, 0x4A)    # support services
SUPPORT_BORDER = ACCENT_GREEN
PROVIDER_COLOR = RGBColor(0x6B, 0x1A, 0x1A)   # external providers
PROVIDER_BORDER = ACCENT_RED
ADMIN_COLOR = RGBColor(0x55, 0x40, 0x80)       # admin / UI
ADMIN_BORDER = RGBColor(0x99, 0x77, 0xCC)

# ── ROW 1: Main pipeline (y=1.2) ──
# Source Systems
add_box(slide, 0.3, 1.3, 1.7, 1.6,
        "Source\nSystems\n\nOMS\nMagento\nMirakl\nChat\nManual",
        SRC_COLOR, SRC_BORDER, font_size=8)

# Arrow: Sources -> Event Ingestion
add_arrow(slide, 2.0, 2.1, 2.5, 2.1)
add_arrow_head(slide, 2.5, 2.1, "right")
add_text_box(slide, 2.05, 1.7, 0.6, 0.3, "REST /\nAMQP", font_size=7, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Event Ingestion Service
add_box(slide, 2.65, 1.5, 1.7, 1.2,
        "Event Ingestion\nService\n:3151\nNormalize & Validate",
        CORE_COLOR, CORE_BORDER, font_size=9)

# Arrow: EIS -> RabbitMQ
add_arrow(slide, 4.35, 2.1, 4.85, 2.1)
add_arrow_head(slide, 4.85, 2.1, "right")

# RabbitMQ (central)
add_box(slide, 4.95, 1.3, 1.8, 1.6,
        "RabbitMQ\n\n6 Exchanges\n23 Queues\n24 Bindings",
        MQ_COLOR, MQ_BORDER, font_size=9)

# Arrow: RabbitMQ -> Notification Engine
add_arrow(slide, 6.75, 2.1, 7.25, 2.1)
add_arrow_head(slide, 7.25, 2.1, "right")

# Notification Engine Service
add_box(slide, 7.35, 1.5, 1.7, 1.2,
        "Notification\nEngine\n:3152\nRules & Orchestrate",
        CORE_COLOR, CORE_BORDER, font_size=9)

# Arrow: Engine -> Channel Router
add_arrow(slide, 9.05, 2.1, 9.55, 2.1)
add_arrow_head(slide, 9.55, 2.1, "right")

# Channel Router Service
add_box(slide, 9.65, 1.5, 1.7, 1.2,
        "Channel Router\nService\n:3154\nRoute & Deliver",
        CORE_COLOR, CORE_BORDER, font_size=9)

# Arrow: Channel Router -> Provider Adapters
add_arrow(slide, 11.35, 2.1, 11.85, 2.1)
add_arrow_head(slide, 11.85, 2.1, "right")

# Provider Adapters
add_box(slide, 11.95, 1.3, 1.1, 1.6,
        "Provider\nAdapters\n\nMailgun\nWhatsApp\nBraze\nAWS SES",
        PROVIDER_COLOR, PROVIDER_BORDER, font_size=8)

# ── Template Service (below Engine, connected) ──
add_box(slide, 7.35, 3.2, 1.7, 0.9,
        "Template Service\n:3153\nRender (Handlebars)",
        CORE_COLOR, CORE_BORDER, font_size=9)

# Arrow: Engine <-> Template (vertical, bidirectional)
add_arrow(slide, 8.2, 2.7, 8.2, 3.2)
add_arrow_head(slide, 8.2, 3.2, "down")
add_text_box(slide, 8.4, 2.8, 0.8, 0.3, "sync\nHTTP", font_size=7, color=MED_GRAY)

# ── ROW 2: Support services (y=4.4) ──

# Bulk Upload
add_box(slide, 0.3, 4.4, 1.7, 0.9,
        "Bulk Upload\nService :3158\nXLSX Processing",
        SUPPORT_COLOR, SUPPORT_BORDER, font_size=9)

# Arrow: Bulk Upload -> Event Ingestion (up-right)
add_arrow(slide, 2.0, 4.6, 2.5, 4.6)
add_arrow(slide, 2.5, 2.7, 2.5, 4.6)
add_arrow_head(slide, 2.5, 2.7, "up")
add_text_box(slide, 2.55, 3.5, 0.9, 0.3, "HTTP POST\n/webhooks", font_size=7, color=MED_GRAY)

# Email Ingest
add_box(slide, 0.3, 5.6, 1.7, 0.9,
        "Email Ingest\nService :3157/2525\nSMTP Ingest",
        SRC_COLOR, SRC_BORDER, font_size=9)

# Arrow: Email Ingest -> RabbitMQ (up-right)
add_arrow(slide, 2.0, 6.0, 3.5, 6.0)
add_arrow(slide, 3.5, 2.9, 3.5, 6.0)
add_arrow_head(slide, 3.5, 2.9, "up")
add_text_box(slide, 2.2, 5.55, 1.2, 0.3, "publishes to\nxch.events.incoming", font_size=7, color=MED_GRAY)

# Audit Service (below Channel Router)
add_box(slide, 9.65, 3.2, 1.7, 0.9,
        "Audit Service\n:3156\nTracking & Analytics",
        SUPPORT_COLOR, SUPPORT_BORDER, font_size=9)

# Arrow: multiple services -> Audit (vertical from router)
add_arrow(slide, 10.5, 2.7, 10.5, 3.2)
add_arrow_head(slide, 10.5, 3.2, "down")
add_text_box(slide, 10.7, 2.8, 1.0, 0.3, "fire &\nforget", font_size=7, color=MED_GRAY)

# ── ROW 3: Admin layer (y=4.4 right side) ──

# Admin UI
add_box(slide, 4.5, 4.6, 1.5, 0.9,
        "Admin UI\n:3159\nNext.js",
        ADMIN_COLOR, ADMIN_BORDER, font_size=9)

# Arrow: Admin UI -> Admin Service
add_arrow(slide, 6.0, 5.05, 6.5, 5.05)
add_arrow_head(slide, 6.5, 5.05, "right")

# Admin Service
add_box(slide, 6.6, 4.6, 1.7, 0.9,
        "Admin Service\n:3155\nConfig & Rules",
        ADMIN_COLOR, ADMIN_BORDER, font_size=9)

# Auth RBAC Backend
add_box(slide, 9.0, 4.6, 1.7, 0.9,
        "Auth RBAC\nBackend :3160\nJWT & Users",
        ADMIN_COLOR, ADMIN_BORDER, font_size=9)

# Auth RBAC Frontend
add_box(slide, 9.0, 5.8, 1.7, 0.9,
        "Auth RBAC\nFrontend :3161\nNext.js",
        ADMIN_COLOR, ADMIN_BORDER, font_size=9)

# Arrow: Auth Frontend -> Auth Backend
add_arrow(slide, 9.85, 5.8, 9.85, 5.5)
add_arrow_head(slide, 9.85, 5.5, "up")

# eCommerce Backoffice
add_box(slide, 11.2, 4.6, 1.5, 0.9,
        "eCommerce\nBackoffice\n:3162",
        ADMIN_COLOR, ADMIN_BORDER, font_size=9)

# ── Legend ──
add_shape_rect(slide, 4.0, 6.5, 6.0, 0.8, RGBColor(0x15, 0x15, 0x28), ACCENT_BLUE)
legend_items = [
    (4.15, CORE_COLOR, "Core Pipeline"),
    (5.55, MQ_COLOR, "Message Broker"),
    (6.95, SUPPORT_COLOR, "Support Services"),
    (8.35, ADMIN_COLOR, "Admin / Auth"),
]
for x, color, label in legend_items:
    add_shape_rect(slide, x, 6.7, 0.25, 0.25, color)
    add_text_box(slide, x + 0.3, 6.68, 1.2, 0.3, label, font_size=9, color=LIGHT_GRAY)

# Legend row 2
legend_items2 = [
    (4.15, SRC_COLOR, "Sources / Ingest"),
    (5.55, PROVIDER_COLOR, "External Providers"),
]
for x, color, label in legend_items2:
    add_shape_rect(slide, x, 7.0, 0.25, 0.25, color)
    add_text_box(slide, x + 0.3, 6.98, 1.2, 0.3, label, font_size=9, color=LIGHT_GRAY)


# ─── SLIDE 8: Key Takeaways ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_rect(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)
add_text_box(slide, 0.5, 0.3, 12, 0.7, "Key Takeaways", font_size=32, color=WHITE, bold=True)

takeaways = [
    ("Architecture",
     "14 microservices designed, 8 implemented. Event-driven with RabbitMQ (6 exchanges, 23 queues). "
     "PostgreSQL with schema-per-service pattern. NestJS + TypeScript backend, Next.js frontend."),
    ("Code Quality",
     "~45K lines of production source code, ~59K lines of test code. "
     "Test-to-source ratio averages 1.3x-2.3x across services. "
     "2,375+ unit tests across 243 suites, plus E2E test coverage."),
    ("Development Velocity",
     "Full platform designed in ~10 days (Feb 19-28). "
     "8 services implemented in ~10 days (Feb 23 - Mar 5). "
     "Parallel development streams enabled rapid delivery."),
    ("Next Steps",
     "Priority: admin-service (39 endpoints, config management hub), "
     "auth-rbac-service-backend (19 endpoints, platform security). "
     "Then: auth-rbac-frontend, ecommerce-backoffice, email-ingest-service. "
     "Remaining provider adapters: Braze, AWS SES."),
]

y = 1.3
for i, (title, desc) in enumerate(takeaways):
    colors = [ACCENT_BLUE, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_RED]
    add_shape_rect(slide, 0.5, y, 0.08, 1.1, colors[i])
    add_text_box(slide, 0.8, y, 11.5, 0.35, title, font_size=18, color=colors[i], bold=True)
    add_text_box(slide, 0.8, y + 0.35, 11.5, 0.8, desc, font_size=13, color=LIGHT_GRAY)
    y += 1.4


# ─── SLIDE 9: Vibe Coding vs Spec-Driven Development ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_rect(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)
add_text_box(slide, 0.5, 0.2, 12, 0.6, "Vibe Coding vs Spec-Driven Development",
             font_size=32, color=WHITE, bold=True)
add_text_box(slide, 0.5, 0.75, 12, 0.3,
             "Two AI-assisted development approaches — same tools, fundamentally different outcomes",
             font_size=14, color=LIGHT_GRAY)

VIBE_COLOR = RGBColor(0xFF, 0x6B, 0x00)
VIBE_BG = RGBColor(0x3A, 0x28, 0x15)
SPEC_COLOR = ACCENT_BLUE
SPEC_BG = RGBColor(0x15, 0x28, 0x3A)

# ── Left card: Vibe Coding ──
add_shape_rect(slide, 0.4, 1.3, 5.8, 2.5, VIBE_BG, VIBE_COLOR)
add_text_box(slide, 0.6, 1.4, 5.4, 0.4, "Vibe Coding", font_size=24, color=VIBE_COLOR, bold=True)
add_text_box(slide, 0.6, 1.85, 5.4, 0.3, '"Prompt and see what happens"',
             font_size=13, color=LIGHT_GRAY)
add_text_box(slide, 0.6, 2.3, 5.4, 1.3,
             "Loose prompts, no upfront design, iterate by feel.\n"
             "Fast for prototypes — but inconsistent quality,\n"
             "no test strategy, hard to maintain past ~1K LoC.\n"
             "Best for: hackathons, MVPs, throwaway projects.",
             font_size=12, color=LIGHT_GRAY)

# ── Right card: Spec-Driven ──
add_shape_rect(slide, 7.1, 1.3, 5.8, 2.5, SPEC_BG, SPEC_COLOR)
add_text_box(slide, 7.3, 1.4, 5.4, 0.4, "Spec-Driven Development", font_size=24, color=SPEC_COLOR, bold=True)
add_text_box(slide, 7.3, 1.85, 5.4, 0.3, '"Design first, then instruct the AI precisely"',
             font_size=13, color=LIGHT_GRAY)
add_text_box(slide, 7.3, 2.3, 5.4, 1.3,
             "Detailed design docs, DB schemas, API contracts —\n"
             "all defined before writing line 1 of code.\n"
             "AI executes against precise specifications.\n"
             "Best for: production systems, teams, long-lived projects.",
             font_size=12, color=LIGHT_GRAY)

# ── Why Spec-Driven Matters ──
add_text_box(slide, 0.5, 4.2, 12, 0.5, "Why Spec-Driven Development Matters",
             font_size=24, color=ACCENT_GREEN, bold=True)

reasons = [
    ("AI is only as good as its instructions.",
     "Vague prompts produce vague code. Detailed specs give the AI exact constraints — data models, error codes, "
     "API contracts, queue topologies — so it generates production-grade code on the first pass, not the fifth."),
    ("Specs are the single source of truth for the entire team.",
     "When every service has a design doc, DB script, and endpoint reference, any developer (or AI) can pick up "
     "any service and know exactly what it should do. No tribal knowledge, no guessing."),
    ("You can scale without chaos.",
     "This project: 14 microservices, 45K source LoC, 59K test LoC, 2,375+ unit tests — designed and built "
     "in ~15 days. That velocity is only possible because specs eliminated ambiguity and rework."),
    ("Specs survive the code.",
     "Code gets refactored, rewritten, or replaced. Specs document the WHY and the WHAT — they remain valuable "
     "long after the current implementation changes. They are the institutional memory of the system."),
]

y = 4.8
for i, (title, desc) in enumerate(reasons):
    add_shape_rect(slide, 0.5, y, 0.08, 0.85, ACCENT_GREEN)
    add_text_box(slide, 0.8, y, 12, 0.3, title, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, 0.8, y + 0.3, 12, 0.55, desc, font_size=11, color=LIGHT_GRAY)
    y += 0.95


# Save
output_path = r"C:\Users\omar.ola\NotificationAPI\temporary\Notification_API_Microservices_Overview.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
